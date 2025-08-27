"""Tests for Office document (Word/PowerPoint) file handling."""

import sys
from io import BytesIO
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from docx import Document
from pptx import Presentation

# Add project root to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from app.file_handler import (
    SUPPORTED_OFFICE_EXTENSIONS,
    extract_text_from_docx,
    extract_text_from_file,
    extract_text_from_pptx,
    process_uploaded_file_with_gemini,
    validate_file_type,
)


class TestOfficeFileTypeValidation:
    """Test cases for Office file type validation."""

    def test_validate_docx_file_type(self):
        """Test that .docx files are recognized as valid."""
        assert validate_file_type("document.docx") is True
        assert validate_file_type("Document.DOCX") is True  # Case insensitive

    def test_validate_pptx_file_type(self):
        """Test that .pptx files are recognized as valid."""
        assert validate_file_type("presentation.pptx") is True
        assert validate_file_type("Presentation.PPTX") is True  # Case insensitive

    def test_office_extensions_defined(self):
        """Test that Office extensions are properly defined."""
        assert ".docx" in SUPPORTED_OFFICE_EXTENSIONS
        assert ".pptx" in SUPPORTED_OFFICE_EXTENSIONS


class TestWordDocumentExtraction:
    """Test cases for Word document text extraction."""

    def test_extract_text_from_docx_with_paragraphs(self):
        """Test extracting text from a Word document with paragraphs."""
        # Create a simple Word document in memory
        doc = Document()
        doc.add_paragraph("Title of Document")
        doc.add_paragraph("This is the first paragraph.")
        doc.add_paragraph("This is the second paragraph.")

        # Save to BytesIO
        doc_bytes = BytesIO()
        doc.save(doc_bytes)
        doc_bytes.seek(0)

        # Extract text
        result = extract_text_from_docx(doc_bytes.read())

        assert "Title of Document" in result
        assert "This is the first paragraph." in result
        assert "This is the second paragraph." in result

    def test_extract_text_from_docx_with_table(self):
        """Test extracting text from a Word document with a table."""
        # Create a Word document with a table
        doc = Document()
        doc.add_paragraph("Document with Table")

        # Add a table
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(0, 2).text = "Header 3"
        table.cell(1, 0).text = "Data 1"
        table.cell(1, 1).text = "Data 2"
        table.cell(1, 2).text = "Data 3"

        # Save to BytesIO
        doc_bytes = BytesIO()
        doc.save(doc_bytes)
        doc_bytes.seek(0)

        # Extract text
        result = extract_text_from_docx(doc_bytes.read())

        assert "Document with Table" in result
        assert "Header 1" in result
        assert "Data 3" in result
        assert "|" in result  # Table cells are joined with |

    def test_extract_text_from_empty_docx(self):
        """Test extracting text from an empty Word document."""
        # Create an empty Word document
        doc = Document()

        # Save to BytesIO
        doc_bytes = BytesIO()
        doc.save(doc_bytes)
        doc_bytes.seek(0)

        # Extract text
        result = extract_text_from_docx(doc_bytes.read())

        assert result == "Word文書にテキストが見つかりませんでした"

    @patch("app.file_handler.logger")
    def test_extract_text_from_invalid_docx(self, mock_logger):
        """Test handling of invalid Word document bytes."""
        invalid_bytes = b"This is not a valid docx file"

        result = extract_text_from_docx(invalid_bytes)

        assert result == "Word文書の読み取りに失敗しました"
        mock_logger.error.assert_called_once()


class TestPowerPointExtraction:
    """Test cases for PowerPoint presentation text extraction."""

    def test_extract_text_from_pptx_with_slides(self):
        """Test extracting text from a PowerPoint with multiple slides."""
        # Create a PowerPoint presentation
        prs = Presentation()

        # Add first slide with title and content
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Presentation Title"
        slide1.placeholders[1].text = "Subtitle text"

        # Add second slide with bullet points
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Slide 2 Title"
        slide2.placeholders[1].text = "• Point 1\n• Point 2\n• Point 3"

        # Save to BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        # Extract text
        result = extract_text_from_pptx(pptx_bytes.read())

        assert "[スライド 1]" in result
        assert "Presentation Title" in result
        assert "Subtitle text" in result
        assert "[スライド 2]" in result
        assert "Slide 2 Title" in result
        assert "Point 1" in result

    def test_extract_text_from_pptx_with_notes(self):
        """Test extracting text including slide notes."""
        # Create a PowerPoint presentation
        prs = Presentation()

        # Add slide with notes
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide with Notes"

        # Add notes to the slide
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "These are speaker notes for the presentation."

        # Save to BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        # Extract text
        result = extract_text_from_pptx(pptx_bytes.read())

        assert "Slide with Notes" in result
        assert "[ノート]: These are speaker notes" in result

    def test_extract_text_from_empty_pptx(self):
        """Test extracting text from an empty PowerPoint."""
        # Create an empty presentation
        prs = Presentation()

        # Save to BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        # Extract text
        result = extract_text_from_pptx(pptx_bytes.read())

        assert result == "PowerPoint文書にテキストが見つかりませんでした"

    @patch("app.file_handler.logger")
    def test_extract_text_from_invalid_pptx(self, mock_logger):
        """Test handling of invalid PowerPoint bytes."""
        invalid_bytes = b"This is not a valid pptx file"

        result = extract_text_from_pptx(invalid_bytes)

        assert result == "PowerPoint文書の読み取りに失敗しました"
        mock_logger.error.assert_called_once()


class TestOfficeFileIntegration:
    """Integration tests for Office file processing."""

    def test_extract_text_from_file_docx(self):
        """Test that extract_text_from_file correctly handles .docx files."""
        # Create a simple Word document
        doc = Document()
        doc.add_paragraph("Test Word Document")

        # Save to BytesIO
        doc_bytes = BytesIO()
        doc.save(doc_bytes)
        doc_bytes.seek(0)

        # Extract using the main function
        result = extract_text_from_file(doc_bytes.read(), "test.docx")

        assert "Test Word Document" in result

    def test_extract_text_from_file_pptx(self):
        """Test that extract_text_from_file correctly handles .pptx files."""
        # Create a simple PowerPoint
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test PowerPoint"

        # Save to BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        # Extract using the main function
        result = extract_text_from_file(pptx_bytes.read(), "test.pptx")

        assert "Test PowerPoint" in result

    @patch("app.file_handler.validate_file_size")
    @patch("app.file_handler.validate_file_type")
    def test_process_uploaded_file_with_gemini_docx(self, mock_validate_type, mock_validate_size):
        """Test processing Word document with process_uploaded_file_with_gemini."""
        mock_validate_size.return_value = True
        mock_validate_type.return_value = True

        # Create a Word document
        doc = Document()
        doc.add_paragraph("Patent specification draft")

        # Save to BytesIO
        doc_bytes = BytesIO()
        doc.save(doc_bytes)
        doc_content = doc_bytes.getvalue()

        # Create mock uploaded file
        file_mock = MagicMock()
        file_mock.name = "patent.docx"
        file_mock.type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        file_mock.size = len(doc_content)
        file_mock.read.return_value = doc_content

        # Process file
        result = process_uploaded_file_with_gemini(file_mock, "Patent document")

        # Verify results
        assert result["filename"] == "patent.docx"
        assert result["comment"] == "Patent document"
        assert "Patent specification draft" in result["extracted_text"]
        assert result["gemini_file_id"] is None  # Office docs use local processing
        assert result["gemini_mime_type"] is None

    @patch("app.file_handler.validate_file_size")
    @patch("app.file_handler.validate_file_type")
    def test_process_uploaded_file_with_gemini_pptx(self, mock_validate_type, mock_validate_size):
        """Test processing PowerPoint with process_uploaded_file_with_gemini."""
        mock_validate_size.return_value = True
        mock_validate_type.return_value = True

        # Create a PowerPoint presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Innovation Presentation"

        # Save to BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_content = pptx_bytes.getvalue()

        # Create mock uploaded file
        file_mock = MagicMock()
        file_mock.name = "innovation.pptx"
        file_mock.type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        file_mock.size = len(pptx_content)
        file_mock.read.return_value = pptx_content

        # Process file
        result = process_uploaded_file_with_gemini(file_mock, "Innovation slides")

        # Verify results
        assert result["filename"] == "innovation.pptx"
        assert result["comment"] == "Innovation slides"
        assert "Innovation Presentation" in result["extracted_text"]
        assert result["gemini_file_id"] is None  # Office docs use local processing
        assert result["gemini_mime_type"] is None


class TestTextLengthLimiting:
    """Test cases for text length limiting in Office documents."""

    def test_docx_text_truncation(self):
        """Test that long Word document text is truncated properly."""
        # Create a Word document with lots of text
        doc = Document()
        for i in range(200):
            doc.add_paragraph(f"This is paragraph number {i} with some text content.")

        # Save to BytesIO
        doc_bytes = BytesIO()
        doc.save(doc_bytes)
        doc_bytes.seek(0)

        # Extract text
        result = extract_text_from_docx(doc_bytes.read())

        # Check truncation
        assert len(result) <= 4000 + 50  # MAX_TEXT_LENGTH plus suffix
        assert "... (以下省略)" in result

    def test_pptx_text_truncation(self):
        """Test that long PowerPoint text is truncated properly."""
        # Create a PowerPoint with many slides
        prs = Presentation()
        for i in range(100):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i} Title"
            slide.placeholders[
                1
            ].text = f"This is the content for slide {i} with detailed information."

        # Save to BytesIO
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        # Extract text
        result = extract_text_from_pptx(pptx_bytes.read())

        # Check truncation
        assert len(result) <= 4000 + 50  # MAX_TEXT_LENGTH plus suffix
        assert "... (以下省略)" in result


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
